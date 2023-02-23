import pandas as pd
import numpy as np
from pptx import slide
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LEGEND_POSITION
from typing import Union

def get_percentages(df: pd.Series) -> pd.DataFrame:
    """Gets the percentages of each unique value in the series

    Args:
        df (pd.Series): A pandas series

    Returns:
        pd.DataFrame: A dataframe of the percentages of each unique value
        
    Examples
    --------
    >>> import pandas as pd
    >>> import excel_charts as e_chart
    >>> data = [
        ["c", 2],
        [4, 5],
        ["c", 5],
        ["c", 4]
    ]
    >>> df = pd.DataFrame(data=data, columns=['A', 'B'])
    >>> df
       A  B
    0  c  2
    1  4  5
    2  c  5
    3  c  4
    >>> df_percentages_A = e_chart.get_percentages(df['A'])
    >>> df_percentages_A
    Name  Percentages
    0    4         25.0
    1    c         75.0
    >>> df_percentages_B = e_chart.get_percentages(df['B'])
    >>> df_percentages_B
        Name  Percentages
    0    2         25.0
    1    4         25.0
    2    5         50.0
    """
    uniques, counts = np.unique(df.astype(str), return_counts=True)
    percentages = dict(zip(uniques, (counts * 100 / len(df))))
    df_percent = pd.DataFrame(data=list(percentages.keys()), columns=['Name'])
    df_percent['Percentages'] = percentages.values()
    return df_percent.round(0)

class H_Bar:
    """
    Creates a new Horizontal Excel Chart Bar Chart instance
    """
    def __init__(self) -> None:
        pass

    def add_chart_data(self, title:str, categories:pd.Series, data:pd.Series) -> None:
        """Adds data for excel chart

        Args:
            title (str): The title of the created chart. default is Untitled
            categories (pd.Series): The categories for your chart
            data (pd.Series): The data that your chart will use
        """
        
        # create actual chart
        self.__chart_data__ = CategoryChartData()
        self.__chart_data__.categories = categories
        self.__chart_data__.add_series(title, data)
        
    def add_chart(self, slide:slide.Slide, location:dict[str,int]={'x':1, 'y':1, 'width':8, 'height':6}) -> None:
        """Creates and adds an excel chart to your slide

        Args:
            slide (slide.Slide): The slide that the chart will be added to.
            location (_type_, optional): The position of the chart. Defaults to {'x':1, 'y':1, 'end_x':8, 'end_y':6}.
        """
        self.__slide__ = slide
        x, y, cx, cy = Inches(location['x']), Inches(location['y']), Inches(location['width']), Inches(location['height'])
        self.__graphic_frame__ = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, self.__chart_data__
        )
    
    @staticmethod
    def __data_labels__(plot, **kwargs) -> None:
        if(kwargs.get('data_labels')):
            if(kwargs.get('label_color')):
                kwargs['label_color'] = H_Bar.__clean_color__(kwargs['label_color'])
        
            plot.has_data_labels = kwargs['data_labels']
            data_labels = plot.data_labels

            # check for label size
            if(kwargs.get('label_size')):
                data_labels.font.size = Pt(kwargs['label_size'])
            else:
                data_labels.font.size = Pt(12)
            
            if(kwargs.get('font')):
                data_labels.font.name = kwargs.get('font')
            else:
                data_labels.font.name = 'Arial'
            
            
            # check for label color
            if(kwargs.get('label_color')):
                data_labels.font.color.rgb = RGBColor.from_string(kwargs['label_color'])
            else:
                data_labels.font.color.rgb = RGBColor.from_string('000000')
        else:
            plot.has_data_labels = False

    @staticmethod
    def __gridlines__(value_axis, **kwargs) -> None:
        if(kwargs.get('major_gridlines')):
            value_axis.has_major_gridlines = kwargs['major_gridlines']
        else:
            value_axis.has_major_gridlines = False
        
        if(kwargs.get('minor_gridlines')):
            value_axis.has_minor_gridlines = kwargs['minor_gridlines']
        else:
            value_axis.has_minor_gridlines = False
            
    @staticmethod
    def __setall__() -> dict[str,Union[int, bool, str, list[str]]]:
        return {
                'axis_labels':True,
                'major_gridlines':True,
                'minor_gridlines':True,
                'data_labels':True,
                'label_color':'#000000',
                'label_size':12,
                'font':'Arial',
                'legend':'right',
                'auto_center':True
        }
       
    @staticmethod
    def __clean_color__(color:str) -> str:
        if(color[0] == '#'):
            return color[1:]
        return color

    def __set_color__(self, chart, **kwargs):
        colors = ['ff0000', '00ff00', '0000ff']

        if(kwargs.get('chart_colors')):
            # set the color of each data point in the first series
            for index in range(len(kwargs['chart_colors'])):
                kwargs['chart_colors'][index] = H_Bar.__clean_color__(kwargs['chart_colors'][index])
            colors = kwargs['chart_colors']

        for idx, point in enumerate(chart.series[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(colors[idx]) # type: ignore

    def __legend__(self, chart, **kwargs) -> None:
        if(kwargs.get('legend')):
            chart.has_legend = True
            
            if(kwargs['legend'] == 'left'):
               chart.legend.position = XL_LEGEND_POSITION.LEFT
            if(kwargs['legend'] == 'right'):
               chart.legend.position = XL_LEGEND_POSITION.RIGHT
            if(kwargs['legend'] == 'top'):
               chart.legend.position = XL_LEGEND_POSITION.TOP
            if(kwargs['legend'] == 'bottom'):
               chart.legend.position = XL_LEGEND_POSITION.BOTTOM
               
            chart.legend.include_in_layout = False
        
            if(kwargs.get('legend_font_size')):
                chart.legend.font.size = Pt(int(kwargs['legend_font_size']))
            else:
                chart.legend.font.size = Pt(18)

    def set_attributes(self, **kwargs:Union[bool,str,int, list[str]]) -> None:
        """
        Adds or removes attributes from excel chart


        Kwargs:
        --------
            axis_labels (bool): Shows or hides chart axis labels. default is False.
            major_gridlines (bool): Shows or hides chart major gridlines. default is False.
            minor_gridlines (bool): Shows or hides chart minor gridlines. default is False.
            data_labels (bool): Shows or hides chart data labels. default is False.
            label_color (hex decimal string): Color for data labels. Only applies if data_labels is set to true. default is black. 
            label_size (int): Size of data_labels, will be measured in points. default is 12pt.
            font (str): type of font for the data labels. default is Arial.
            chart_colors (list[hex decimal]): List of hex decimal strings for the chart colors. default is ['ff0000', '00ff00', '0000ff'].
            legend (str): Adds a legend to the chart in the specified position. default is None. Options are (left, right, top, bottom).
            legend_font_size (int): Side of the font in the legend. default is 18. Only applicable if legend is on.
            auto_center (bool): Centers the chart on the slide. default is True.
            all (bool): Enables all attributes. will override all other settings
        """
        if(kwargs.get('all')):
            kwargs = self.__setall__()
        
        chart = self.__graphic_frame__.chart
        
        # adds legend to the chart
        self.__legend__(chart, **kwargs)

        # set bar colors
        self.__set_color__(chart, **kwargs)
        
        if(kwargs.get('axis_labels')):
            chart.value_axis.visible = kwargs['axis_labels']
        else:
            chart.value_axis.visible = False

        # Plot
        plot = chart.plots[0]

        # create data_labels
        self.__data_labels__(plot, **kwargs)

        # Value axis
        value_axis = chart.value_axis
        self.__gridlines__(value_axis, **kwargs)